# -*- coding: utf-8 -*-
"""Build the analysis report deck (16:9) with a FIXED grid, consistent title/core-message/
bullet baselines, card layouts, page numbers — per the user's alignment spec.
All content is drawn from analysis/I1_REPORT.md etc. (algorithmic results)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ---------------- palette & grid ----------------
NAVY   = RGBColor(8, 31, 70)      # 제목 기본
DGRAY  = RGBColor(80, 80, 80)     # 핵심메시지/캡션
GREEN  = RGBColor(0, 150, 95)     # accent
PGRAY  = RGBColor(120, 120, 120)  # 페이지번호
WHITE  = RGBColor(255, 255, 255)
LIGHT  = RGBColor(238, 242, 248)  # 카드 배경(연한 네이비톤)
LIGHTG = RGBColor(232, 245, 240)  # 카드 배경(연한 그린톤)
LINE   = RGBColor(208, 216, 228)
BODYTX = RGBColor(40, 45, 55)
FONT   = "맑은 고딕"

# slide 16:9
SW, SH = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.62)
CONTENT_W = SW - 2*MARGIN
# fixed vertical grid
TITLE_TOP, TITLE_H = Inches(0.40), Inches(0.74)     # ~12-15%
RULE_TOP           = Inches(1.16)                    # accent rule under title
CORE_TOP, CORE_H   = Inches(1.28), Inches(0.56)     # ~8-10%
BODY_TOP, BODY_H   = Inches(2.00), Inches(4.95)      # ~66%
FOOT_TOP           = Inches(7.04)

prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]

# ---------------- low-level helpers ----------------
def _set_font(run, size, bold, color, italic=False):
    run.font.name = FONT; run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color
    # ensure east-asian font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
    ea.set('typeface', FONT)

def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    return tb, tf

def para(tf, text, size, bold=False, color=BODYTX, align=PP_ALIGN.LEFT,
         space_after=6, level=0, first=False, italic=False, bullet=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.level = level
    p.space_after = Pt(space_after); p.space_before = Pt(0); p.line_spacing = 1.04
    prefix = ("• " if bullet=="•" else ("– " if bullet=="–" else ""))
    r = p.add_run(); r.text = prefix + text
    _set_font(r, size, bold, color, italic)
    return p

def rrect(slide, l, t, w, h, fill, line=None, line_w=0.75, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = slide.shapes.add_shape(shape, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp

# ---------------- structural helpers ----------------
def base(title, core=None, core_color=DGRAY, page=True, tag=None):
    s = prs.slides.add_slide(BLANK)
    # optional section tag (small, above title baseline-right)
    if tag:
        tb,tf = textbox(s, MARGIN, Inches(0.12), CONTENT_W, Inches(0.26))
        para(tf, tag, 11, bold=True, color=GREEN, first=True, space_after=0)
    # title
    tb, tf = textbox(s, MARGIN, TITLE_TOP, CONTENT_W, TITLE_H, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, title, 26, bold=True, color=NAVY, first=True, space_after=0)
    # accent rule
    ln = rrect(s, MARGIN, RULE_TOP, CONTENT_W, Pt(2.2), NAVY, shape=MSO_SHAPE.RECTANGLE)
    # core message
    if core:
        tb, tf = textbox(s, MARGIN, CORE_TOP, CONTENT_W, CORE_H, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, core, 15, bold=True, color=core_color, first=True, space_after=0)
    # page number
    if page:
        n = len(prs.slides.__iter__.__self__._sldIdLst)  # current count
        tb, tf = textbox(s, SW-Inches(1.2), FOOT_TOP, Inches(0.9), Inches(0.3))
        para(tf, str(len(prs.slides._sldIdLst)), 10, color=PGRAY, align=PP_ALIGN.RIGHT, first=True, space_after=0)
    return s

def caption(slide, text, l, t, w):
    tb, tf = textbox(slide, l, t, w, Inches(0.26))
    para(tf, text, 9.5, color=DGRAY, first=True, space_after=0)

def bullets(slide, items, l=MARGIN, t=BODY_TOP, w=CONTENT_W, h=BODY_H, size=15, gap=9, anchor=MSO_ANCHOR.TOP):
    tb, tf = textbox(slide, l, t, w, h, anchor=anchor)
    for i,(txt, lv) in enumerate(items):
        para(tf, txt, size if lv==0 else size-2, bold=(lv==0 and txt.endswith(':')),
             color=BODYTX if lv==0 else DGRAY, first=(i==0), space_after=gap,
             level=lv, bullet="•" if lv==0 else "–")
    return tb

def table(slide, data, l, t, w, h, col_w=None, header=True, fs=11.5, hfs=11.5):
    rows, cols = len(data), len(data[0])
    gtbl = slide.shapes.add_table(rows, cols, l, t, w, h).table
    gtbl.first_row = header; gtbl.horz_banding = True
    if col_w:
        tot = sum(col_w)
        for j,cw in enumerate(col_w):
            gtbl.columns[j].width = Emu(int(w * cw/tot))
    for i,row in enumerate(data):
        for j,val in enumerate(row):
            c = gtbl.cell(i,j); c.margin_left=Pt(5); c.margin_right=Pt(5)
            c.margin_top=Pt(2); c.margin_bottom=Pt(2)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf=c.text_frame; tf.word_wrap=True
            p=tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER
            r=p.add_run(); r.text=str(val)
            if i==0 and header:
                c.fill.solid(); c.fill.fore_color.rgb=NAVY
                _set_font(r, hfs, True, WHITE)
            else:
                c.fill.solid(); c.fill.fore_color.rgb = WHITE if i%2 else LIGHT
                _set_font(r, fs, False, BODYTX)
    return gtbl

def card(slide, l, t, w, h, title, lines, accent=NAVY, bg=LIGHT):
    rrect(slide, l, t, w, h, bg, line=LINE, line_w=0.75)
    rrect(slide, l, t, w, Inches(0.40), accent)  # header strip
    tb, tf = textbox(slide, l+Inches(0.14), t+Inches(0.02), w-Inches(0.24), Inches(0.38), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, title, 13, bold=True, color=WHITE, first=True, space_after=0)
    tb, tf = textbox(slide, l+Inches(0.14), t+Inches(0.48), w-Inches(0.26), h-Inches(0.56))
    for i,ln in enumerate(lines):
        para(tf, ln, 11.5, color=BODYTX, first=(i==0), space_after=5, bullet="•")

def quad(slide, blocks):
    """2x2 issue layout: blocks = [(title,lines,accent), x4] in order
       별도처리 / 의미분석 / 대책 / 추가분석."""
    gx, gy = Inches(0.32), Inches(0.30)
    cw = (CONTENT_W - gx)/2; ch = (BODY_H - gy)/2
    pos = [(MARGIN, BODY_TOP), (MARGIN+cw+gx, BODY_TOP),
           (MARGIN, BODY_TOP+ch+gy), (MARGIN+cw+gx, BODY_TOP+ch+gy)]
    for (l,t),(title,lines,accent,bg) in zip(pos, blocks):
        card(slide, l, t, cw, ch, title, lines, accent, bg)

def chev_flow(slide, steps, t, h=Inches(0.95)):
    n=len(steps); gap=Inches(0.10)
    w=(CONTENT_W-gap*(n-1))/n
    for i,(label,sub) in enumerate(steps):
        l=MARGIN+i*(w+gap)
        sp=rrect(slide, l, t, w, h, NAVY if i%2==0 else GREEN, shape=MSO_SHAPE.PENTAGON)
        tf=sp.text_frame; tf.word_wrap=True
        p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=label; _set_font(r,12.5,True,WHITE)
        p2=tf.add_paragraph(); p2.alignment=PP_ALIGN.CENTER
        r2=p2.add_run(); r2.text=sub; _set_font(r2,9.5,False,WHITE)

# =================================================================
# SLIDE 1 — 표지
# =================================================================
s = prs.slides.add_slide(BLANK)
rrect(s, 0, 0, SW, SH, NAVY, shape=MSO_SHAPE.RECTANGLE)
rrect(s, 0, Inches(5.0), SW, Pt(3), GREEN, shape=MSO_SHAPE.RECTANGLE)
tb,tf = textbox(s, MARGIN, Inches(2.5), SW-2*MARGIN, Inches(2.0), anchor=MSO_ANCHOR.MIDDLE)
para(tf, "공공데이터 ‘데이터 지도’\n거버넌스·인텔리전스 분석 보고서", 34, bold=True, color=WHITE, first=True, space_after=8)
tb,tf = textbox(s, MARGIN, Inches(5.2), SW-2*MARGIN, Inches(1.2))
para(tf, "999개 공공데이터셋(교육·보건·공공행정) 데이터 지도의 비지도 데이터마이닝 분석", 15, color=RGBColor(200,210,225), first=True, space_after=6)
para(tf, "1D 클러스터링 · AOI 롤업 · FP-Growth 빈발패턴 · 연관규칙 · 국소→전체 검증", 13, color=GREEN, space_after=0)

# =================================================================
# SLIDE 2 — 분석 개요 및 목적  (좌 bullet + 우 카드)
# =================================================================
s = base("분석 개요 및 목적", "원본 트랜잭션이 아닌 ‘데이터 지도(메타데이터)’를 비지도 알고리즘으로 분석함", tag="OVERVIEW")
bullets(s, [
    ("대상: 공공데이터포털 999개 데이터셋의 프로파일링 산출물(데이터 지도)", 0),
    ("구성: 컬럼/테이블 프로파일 · 쌍별 값·도메인 중첩 · 유사도 클러스터", 0),
    ("목표: 데이터 거버넌스 이슈와 메타 수준 비즈니스 인텔리전스 도출", 0),
    ("형식: ‘주말>기저귀>맥주’형 연관규칙을 데이터셋·컬럼 수준에서 발굴", 0),
    ("원칙: LLM 추론 배제, 모든 결론은 설명가능한 알고리즘 실행 결과", 0),
], w=Inches(6.7))
# right cards: 5원칙
rx = MARGIN+Inches(6.95); rw = CONTENT_W-Inches(6.95)
card(s, rx, BODY_TOP, rw, Inches(2.30), "분석 5원칙(과제 지침)",
     ["① 비지도 알고리즘·설명가능","② 이산화·AOI·빈발패턴·연관규칙·서브스페이스",
      "③ 데이터 롤업 적극 활용","④ 국소 규칙은 전체에서 검증","⑤ 보조 규칙은 proj 이미지 참조"], accent=GREEN, bg=LIGHTG)
card(s, rx, BODY_TOP+Inches(2.55), rw, Inches(2.10), "도메인 분포(999)",
     ["교육 400건 (40.0%)","보건 300건 (30.0%)","공공행정 299건 (29.9%)","cate1×cate2 깔끔히 분리됨"], accent=NAVY)

# =================================================================
# SLIDE 3 — 대상 데이터 구성 (상단요약 + 표)
# =================================================================
s = base("대상 데이터 구성", "9개 파일 = 프로파일 + 쌍별 중첩 + 유사도 클러스터의 3계층 구조", tag="DATA")
data=[["파일","역할","규모"],
 ["table_map","데이터셋 메타+분류+통계","999행"],
 ["column_map","컬럼 프로파일(키품질·고유값 등)","9,840행"],
 ["distinct_value","컬럼별 distinct 값+빈도","803만행/635MB"],
 ["table/column_summary_map","객체별 중첩 통계 롤업","996/8,360행"],
 ["combined_pair_map","컬럼쌍 조인가능성(linking_ratio)","74만행/300MB"],
 ["table_pair_map","테이블쌍 유사도","25만행/296MB"],
 ["column_pair_map","컬럼쌍 값·도메인 중첩·중복","524만행/3.4GB"],
 ["*_value_tree / *_str_tree","값·이름 유사도 계층 클러스터","331~4,549행"]]
table(s, data, MARGIN, BODY_TOP, CONTENT_W, Inches(4.4), col_w=[2.2,4.2,1.6], fs=12, hfs=12.5)
caption(s, "출처: data/ 폴더 CSV(UTF-8). 규모는 scripts/05 행수 집계 결과", MARGIN, Inches(6.55), CONTENT_W)

# =================================================================
# SLIDE 4 — 분석 방법론 (파이프라인 도식 + 알고리즘 매핑)
# =================================================================
s = base("분석 방법론 파이프라인", "이산화 → AOI 롤업 → 빈발패턴 → 연관규칙 → 국소·전체 검증의 5단계", tag="METHOD")
chev_flow(s, [("① 이산화","1D KMeans"),("② AOI 롤업","개념계층"),
              ("③ 빈발패턴","FP-Growth"),("④ 연관규칙","lift 등"),("⑤ 검증","규칙#4")],
          t=BODY_TOP+Inches(0.1))
data=[["단계","알고리즘","적용"],
 ["이산화","log1p + 1D KMeans","rec_cnt·column_cnt·수치비 → 순서형 구간"],
 ["AOI 롤업","정규식 개념계층","table_source→기관유형, col_nm→컬럼개념"],
 ["빈발패턴","FP-Growth(max_len=3)","테이블=거래, 속성·컬럼=아이템"],
 ["연관규칙","association_rules(lift)","support·confidence·lift·leverage·conviction"],
 ["검증","도메인별→전체 재계산","국소 강규칙의 일반화 여부 판정"]]
table(s, data, MARGIN, BODY_TOP+Inches(1.35), CONTENT_W, Inches(3.0), col_w=[1.4,2.6,4.6], fs=11.5, hfs=12)
caption(s, "재현: scripts/06_i1_assoc_rules.py · 상세: analysis/I1_REPORT.md §0.1", MARGIN, Inches(6.6), CONTENT_W)

# =================================================================
# SLIDES 5-8 — 공통 데이터 준비
# =================================================================
# 5: 적재/인코딩/형변환
s = base("공통 데이터 준비 ① 적재·인코딩·형변환", "모든 이슈가 공유하는 기반 전처리 — 인코딩 확정과 결정적 형변환", core_color=DGRAY, tag="공통 전처리")
bullets(s, [
    ("인코딩 확정: 바이트 검증으로 데이터는 UTF-8 정상 — 콘솔(CP949) 표시만 깨짐", 0),
    ("결과물은 항상 UTF-8 파일로 출력해 한글 무결성 확보", 0),
    ("적재: table_map(999) · column_map(9,840)을 dtype=str로 안전 적재 후 형변환", 0),
    ("형변환: 수치 컬럼은 to_numeric(coerce) — 비수치/날짜는 NaN 처리", 0),
    ("정합 기준: table_id 결측 행 제거, 컬럼-테이블 조인키 정합 확인", 0),
])

# 6: 1D 이산화
s = base("공통 데이터 준비 ② 1D 이산화", "멱법칙 수치를 log1p+KMeans로 설명가능한 순서형 구간으로 변환", core_color=DGRAY, tag="공통 전처리")
bullets(s, [
    ("문제: rec_cnt가 1~833,466으로 극단적 멱법칙 → 그대로면 큰 값이 지배", 0),
    ("절차: log1p 변환 → 1차원 KMeans(n_init=10) → 군집중심 정렬로 순서 라벨", 0),
    ("설명가능성: 각 구간을 원척도 [min~max]로 환원해 사람이 읽도록 함", 0),
], w=Inches(6.6))
data=[["변수","구간(원척도, n)"],
 ["레코드수","극소[1-22]·소[23-130]·중[133-829]·대[843-8019]·극대[8578-833466]"],
 ["컬럼폭","협소[1-5]·표준[6-11]·광폭[12-44]·초광폭[49-338]"],
 ["수치비","범주우세[0-.23]·혼합[.23-.58]·수치우세[.60-1.0]"]]
table(s, data, MARGIN, BODY_TOP+Inches(2.55), CONTENT_W, Inches(1.9), col_w=[1.3,6.7], fs=11, hfs=12)
caption(s, "KMeans 경계는 데이터 구동, 라벨은 중심 오름차순 매핑 — scripts/04·06", MARGIN, Inches(6.6), CONTENT_W)

# 7: AOI 롤업
s = base("공통 데이터 준비 ③ AOI 롤업", "999개 자유 기관명과 컬럼명을 상위 개념으로 일반화하여 패턴 가시화", core_color=DGRAY, tag="공통 전처리")
bullets(s, [
    ("출처 일반화: table_source → 정규식 접미사 규칙으로 7개 기관유형", 0),
    ("예) …교육청·대학교→교육기관, …부·청·위원회→중앙행정, …시·군·구→기초지자체", 1),
    ("컬럼명 정규화: 괄호·공백·숫자접두 제거(결정적) → 컬럼개념", 0),
    ("동의어 병합(연번≈순번)은 의도적 미적용 — BI 규칙 오염 방지(별도 G3)", 1),
    ("효과: 빈발패턴이 통계적 의미를 갖도록 차원 축약", 0),
], w=Inches(6.6))
card(s, MARGIN+Inches(6.85), BODY_TOP, CONTENT_W-Inches(6.85), Inches(3.9), "기관유형 분포(AOI 결과)",
     ["기초지자체 404","공공기관 245","미상(소스) 145","광역지자체 83","교육기관 49","중앙행정 39","기타 34"], accent=GREEN, bg=LIGHTG)

# 8: 트랜잭션/FP-Growth/검증
s = base("공통 데이터 준비 ④ 트랜잭션·패턴·검증", "테이블을 거래로 보고 두 종류 아이템으로 FP-Growth·연관규칙·검증 수행", core_color=DGRAY, tag="공통 전처리")
bullets(s, [
    ("트랜잭션: 테이블 1건 = 1 거래(transaction)", 0),
    ("아이템(A) 속성: cate·기관·레코드·컬럼폭·수치비·유일키·상수컬럼", 0),
    ("아이템(B) 장바구니: 테이블이 보유한 컬럼개념 집합(평균 9.6개)", 0),
    ("채굴: TransactionEncoder→FP-Growth→association_rules(lift)", 0),
    ("검증: 도메인별 국소 채굴 후 전체 999건에서 지표 재계산(규칙#4)", 0),
])

# =================================================================
# SLIDE 9 — 발견 종합 (카드 4)
# =================================================================
s = base("분석 결과 종합", "거버넌스 3대 이슈와 인텔리전스 4대 발견을 도출함", core_color=GREEN, tag="SUMMARY")
gx=Inches(0.28); cw=(CONTENT_W-3*gx)/4; t=BODY_TOP; h=Inches(4.6)
card(s, MARGIN, t, cw, h, "거버넌스 G1",
     ["출처 결손의 계통성","‘유아및초·중등교육’ 100건","100% placeholder('소스')","lift 6.89·conf 1.00"], accent=NAVY)
card(s, MARGIN+(cw+gx), t, cw, h, "거버넌스 G2·G3",
     ["죽은 메트릭 3종 전부 0","상수컬럼 1,372개(13.9%)","컬럼명 표준 부재","same_ratio 네이밍 오류"], accent=NAVY)
card(s, MARGIN+2*(cw+gx), t, cw, h, "인텔리전스 I-패턴",
     ["스키마 패밀리=데이터셋 원형","경도→위도 lift 22.7","교과목코드↔분반 124.9","연령대→성별 35.2"], accent=GREEN, bg=LIGHTG)
card(s, MARGIN+3*(cw+gx), t, cw, h, "인텔리전스 I-위험·검증",
     ["의료기관 PII 묶음 lift 46","교차도메인 재식별 위험","규칙#4: 1.00→0.44 약화","카탈로그 오분류 후보"], accent=GREEN, bg=LIGHTG)

# =================================================================
# ISSUE SLIDES (2x2: 별도처리/의미/대책/추가분석)
# =================================================================
def issue(title, core, sep, mean, action, more, core_color=DGRAY, tag="거버넌스 이슈"):
    s = base(title, core, core_color=core_color, tag=tag)
    quad(s, [
        ("① 이슈별 데이터 처리", sep, NAVY, LIGHT),
        ("② 의미 분석", mean, NAVY, LIGHT),
        ("③ 대책", action, GREEN, LIGHTG),
        ("④ 추가 필요 분석", more, GREEN, LIGHTG),
    ])
    return s

# G1
issue("거버넌스 G1 출처 결손의 계통성",
 "출처 누락은 무작위가 아니라 단일 분류에 100% 집중된 계통적 손실임",
 ["AOI로 table_source→기관유형, ‘소스’=미상 식별",
  "cate2별 placeholder 비율 교차표 산출",
  "A1 비자명 연관규칙(lift>1.2)에서 결합 포착"],
 ["placeholder 145건 중 100건이 ‘유아및초·중등교육’(100%)",
  "단일 수집 배치의 provenance 일괄 소실 가능성",
  "유일키·상수컬럼·표준폭 동반(conf 0.79~0.84)=동형 템플릿"],
 ["해당 배치 출처 복구 1순위 지정",
  "적재 파이프라인에 provenance 필수값 검증(누락 차단)",
  "출처 표준 스키마·코드체계 정의"],
 ["regi_date·url로 수집경로 역추적",
  "table_str_tree로 동형 템플릿 군집 확정",
  "타 분류 결손율 추세 모니터링"])

# G2
issue("거버넌스 G2 죽은 메트릭",
 "카탈로그 품질·활용 지표가 전부 미충전되어 분석축으로 쓸 수 없음",
 ["수치 컬럼 분포 통계(%zero/%null) 산출",
  "1D 통계로 상수성·미충전 판별",
  "컬럼 의미 vs 실제 값 정의 대조"],
 ["null_ratio·table_qty_index·cum_search_num=9,840건 전부 0",
  "품질·활용 지표 미계산 또는 미적재",
  "same_ratio는 비율 아님(rec/distinct=평균중복도)→네이밍 오류"],
 ["지표 산출 파이프라인 점검·재적재",
  "컬럼 의미 사전(메타 표준) 수립·명칭 정정",
  "미충전 지표 폐기 또는 채움 정책 결정"],
 ["distinct_value로 null_ratio 실제 재계산",
  "검색로그 연동해 cum_search_num 복원",
  "지표 정의-구현 정합성 감사"])

# G3
issue("거버넌스 G3 컬럼 표준·품질",
 "동일 개념의 상이한 이름과 상수·불일치 컬럼이 표준 부재를 드러냄",
 ["col_nm 빈도 집계(최빈 컬럼)",
  "distinct_cnt·pk_ratio 1D 이산화",
  "column_cnt vs (범주+수치) 정합성 검사"],
 ["데이터기준일자(169)/데이터기준일(46), 연번(168)/순번(79)/번호(34)",
  "상수 컬럼(distinct=1) 1,372개(13.9%)=정보량 0",
  "스키마 회계 불일치 테이블 56개"],
 ["문자유사도 군집으로 표준어 사전 산출(G3)",
  "상수·중복 컬럼 정리 가이드 배포",
  "스키마 검증 룰(컬럼수 정합) 적용"],
 ["편집거리·자모 기반 컬럼명 클러스터링",
  "table_str_tree 활용한 명칭군 검증",
  "표준어 매핑의 도메인 영향 평가"])

# I-스키마 패밀리
issue("인텔리전스 스키마 패밀리",
 "동일 컬럼묶음의 반복 출현이 데이터셋 원형과 통합 후보를 드러냄",
 ["컬럼개념 장바구니 구성(테이블별 보유개념)",
  "FP-Growth 최대빈발 항목집합(≥8테이블·≥3개)",
  "패밀리별 지배 도메인·대표 출처 집계"],
 ["사학재정군·의료기관군·POI군·국가시험군 등 원형 식별",
  "동일 기관이 동일 템플릿 다건 공개(예: 국가시험원 9건)",
  "POI군(경도·위도·주소·전화)은 3개 도메인 횡단 범용"],
 ["POI·의료기관·업소 표준 데이터모델 정의",
  "동일기관 다건의 통합 카탈로깅",
  "패밀리 단위 스키마 표준화"],
 ["table_pair_map 유사도로 near-duplicate 확정",
  "combined_pair_map으로 조인키 검증(I3)",
  "패밀리 커버리지·결측 분석"],
 core_color=GREEN, tag="인텔리전스")

# I-연관규칙
s = base("인텔리전스 흥미로운 연관규칙", "‘주말>기저귀>맥주’ 대응 — 컬럼 동시출현 규칙으로 데이터 구조 지문 도출", core_color=GREEN, tag="인텔리전스")
data=[["주제","규칙","conf","lift"],
 ["지리좌표","경도 → 위도","1.00","22.7"],
 ["기간쌍","교육시작일 → 교육종료일","1.00","90.8"],
 ["행정구역","시도명 → 시군구명","0.94","39.0"],
 ["인구통계","연령대 → 성별","0.85","35.2"],
 ["교육스키마","교과목코드 ↔ 분반","1.00","124.9"],
 ["사학재정","연구비 → {일반·지정기부금}","1.00","111.0"]]
table(s, data, MARGIN, BODY_TOP, Inches(7.0), Inches(3.5), col_w=[1.4,3.6,1.0,1.0], fs=12, hfs=12)
caption(s, "전체 999, lift>2, 실제 테이블 추적 가능 — analysis/I1_REPORT.md §4", MARGIN, Inches(6.4), Inches(7.0))
rx=MARGIN+Inches(7.25); rw=CONTENT_W-Inches(7.25)
card(s, rx, BODY_TOP, rw, Inches(2.0), "읽는 법(예)",
     ["‘경도’ 있으면 거의 100% ‘위도’ 동반","그 동반은 우연 대비 22.7배 흔함","→ 시각화형 POI 데이터셋의 식별 신호"], accent=NAVY)
card(s, rx, BODY_TOP+Inches(2.2), rw, Inches(1.95), "대책·추가분석",
     ["규칙 기반 스키마 자동태깅·무결성 체크","좌표/기간/주소 쌍 결측 검증","max_len↑·저support 규칙 도메인 확인"], accent=GREEN, bg=LIGHTG)

# I-재식별 위험 (G7)
issue("인텔리전스·거버넌스 교차도메인 재식별 위험",
 "명·주소·전화번호가 항상 묶여 공개되어 단일·결합 재식별 위험을 형성함",
 ["의료기관 패밀리+연관규칙(전화→{명,주소} lift 46) 추출",
  "교차도메인 공통 컬럼 집계(전화 140·주소 85·위경도 44)",
  "PII 후보 컬럼 태깅"],
 ["명+주소+전화 묶음=데이터셋 내 직접 재식별 가능",
  "도메인 횡단 공통 식별자로 결합 위험 증폭",
  "양면성: 데이터 연계 기회 ↔ 프라이버시 위험"],
 ["PII 표준 거버넌스(마스킹·집계화) 적용",
  "결합 위험 데이터셋 공개 전 영향평가 의무화",
  "식별자 컬럼 자동 탐지·태깅 체계"],
 ["combined_pair_map linking_ratio로 결합가능성 정량화",
  "k-익명성·재식별 위험도 평가",
  "도메인 횡단 식별자 전수 인벤토리"],
 core_color=GREEN, tag="인텔리전스·거버넌스")

# I-규칙#4 검증
issue("규칙#4 국소→전체 검증",
 "국소에서 강했던 규칙이 전체에서 약화·비대칭됨을 실증함",
 ["도메인별 국소 채굴 후 전체 999건 재계산(verify)",
  "국소·전체 support·conf·lift 병기",
  "방향성(A→B vs B→A) 비대칭 점검"],
 ["소재지도로명→지번: 공공행정 1.00 → 전체 0.44 약화",
  "보건은 지번주소 단독 보유 多 → 방향 붕괴",
  "보건 일부는 데이터 명세서(사전) 파일=오분류 후보"],
 ["규칙 일반화 전 전체검증 의무화",
  "카탈로그 오분류(명세서 vs 실데이터) 재분류",
  "방향성 규칙은 양방향 모두 표기"],
 ["명세서 의심 테이블 별도 군집·식별",
  "비대칭 규칙 전수 점검",
  "도메인별 규칙 안정성 모니터링"],
 core_color=GREEN, tag="검증")

# =================================================================
# SLIDE — 종합 시사점 & 로드맵
# =================================================================
s = base("종합 시사점 및 거버넌스 대책 로드맵", "발견을 즉시·단기·중기 조치로 우선순위화함", core_color=GREEN, tag="ROADMAP")
gx=Inches(0.30); cw=(CONTENT_W-2*gx)/3; t=BODY_TOP; h=Inches(4.5)
card(s, MARGIN, t, cw, h, "즉시(거버넌스)",
     ["출처 결손 배치 provenance 복구","적재 시 필수 메타 검증 도입","죽은 메트릭·네이밍 정정"], accent=NAVY)
card(s, MARGIN+(cw+gx), t, cw, h, "단기(표준화)",
     ["컬럼 표준어 사전(G3) 구축","POI·의료기관 표준 데이터모델","PII 식별자 태깅·마스킹 정책"], accent=NAVY)
card(s, MARGIN+2*(cw+gx), t, cw, h, "중기(인텔리전스)",
     ["조인키 검증으로 데이터 연계(I3)","서브스페이스 클러스터링(I5)","near-duplicate 통합 카탈로깅"], accent=GREEN, bg=LIGHTG)

# =================================================================
# SLIDE — 다음 단계 & 한계
# =================================================================
s = base("다음 단계 분석 및 한계", "심화 분석 경로와 현재 분석의 명시적 한계를 제시함", core_color=DGRAY, tag="NEXT")
bullets(s, [
    ("I3(연계): POI·의료기관 패밀리를 linking_ratio로 실제 조인키 검증", 0),
    ("G3(거버넌스): 문자유사도 군집으로 표준어 사전 산출", 0),
    ("I5(심화): 컬럼 프로파일 서브스페이스 클러스터링으로 원형 정교화", 0),
    ("한계: 장바구니 규칙 support 낮음(8~44 테이블) → 운영 전 도메인 확인 필요", 0),
    ("한계: 활용/품질 지표 미충전, distinct_value 정렬편향, proj 이미지 손상", 0),
])

out = r"d:\data_intelligence_2\report\데이터지도_분석보고서.pptx"
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print("WROTE", out, "slides=", len(prs.slides._sldIdLst))
